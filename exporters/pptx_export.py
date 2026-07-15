"""
exporters/pptx_export.py
---------------------------
PowerPoint (.pptx) exporter — a short slide deck version of the newsletter,
useful for a leadership skim/readout.
"""

from __future__ import annotations
from pathlib import Path
from typing import Dict, Any
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

NAVY = RGBColor(0x1F, 0x4E, 0x78)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


def _title_slide(prs: Presentation, newsletter: Dict[str, Any]):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "FMCG Deal Intelligence Newsletter"
    slide.placeholders[1].text = f"Generated {newsletter.get('generated_at', '')}"
    return slide


def _summary_slide(prs: Presentation, newsletter: Dict[str, Any]):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Executive Summary"
    body = slide.placeholders[1].text_frame
    body.text = newsletter.get("executive_summary", "")
    body.word_wrap = True
    return slide


def _deals_slides(prs: Presentation, newsletter: Dict[str, Any], per_slide: int = 5):
    deals = newsletter.get("deals", [])
    for i in range(0, len(deals), per_slide):
        chunk = deals[i:i + per_slide]
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = f"Top Deals ({i + 1}-{i + len(chunk)})"
        body = slide.placeholders[1].text_frame
        body.word_wrap = True
        first = True
        for d in chunk:
            line = f"{d.get('title', '')[:90]} — {d.get('deal_type', '').replace('_', ' ').title()}"
            if d.get("amount") and d.get("amount") != "—":
                line += f" ({d['amount']})"
            if first:
                body.text = line
                first = False
            else:
                p = body.add_paragraph()
                p.text = line
            p_obj = body.paragraphs[-1]
            p_obj.font.size = Pt(16)


def _trend_slide(prs: Presentation, newsletter: Dict[str, Any]):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Market Trend Snapshot"
    body = slide.placeholders[1].text_frame
    trend = newsletter.get("trend_snapshot", {})
    if not trend:
        body.text = "No trend data available."
        return
    first = True
    for dtype, count in sorted(trend.items(), key=lambda x: -x[1]):
        line = f"{dtype.replace('_', ' ').title()}: {count} deal(s)"
        if first:
            body.text = line
            first = False
        else:
            body.add_paragraph().text = line


def export_pptx(newsletter: Dict[str, Any], out_path: str) -> str:
    Path(out_path).parent.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    _title_slide(prs, newsletter)
    _summary_slide(prs, newsletter)
    _deals_slides(prs, newsletter)
    _trend_slide(prs, newsletter)
    prs.save(out_path)
    return out_path
